from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize the presentation canvas
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 modern widescreen standard
prs.slide_height = Inches(7.5)

# Color Palette: Modern Industrial Slate Theme
BG_COLOR = RGBColor(0x0F, 0x17, 0x2A)       # Sleek Deep Charcoal
TEXT_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)     # Crisp Alpine White
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)     # Cool Steel Gray
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # Emerald Ingestion Success
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # Safety Warning / Action required
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)        # Slate Blue Card Surface

def apply_background(slide):
    """Fills the slide background with deep charcoal."""
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG_COLOR
    bg_shape.line.fill.background()
    return bg_shape

def create_title(slide, text, top=0.5):
    """Injects a standardized top slide header."""
    tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    return tx_box

slide_layout = prs.slide_layouts[6]  # Blank layout canvas

# ==========================================
# SLIDE 1: INTRODUCTION & VISION
# ==========================================
slide1 = prs.slides.add_slide(slide_layout)
apply_background(slide1)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "Project MetaVal: Enterprise Email & PDF Data Ingestion Pipeline"
p1.font.name = "Arial"
p1.font.size = Pt(38)
p1.font.bold = True
p1.font.color.rgb = TEXT_LIGHT
p1.space_after = Pt(14)

p2 = tf.add_paragraph()
p2.text = "Transitioning Manufacturing Operations from Fragmented Manual Logging to Intelligent Automation"
p2.font.name = "Arial"
p2.font.size = Pt(16)
p2.font.color.rgb = ACCENT_GREEN
p2.space_after = Pt(40)

p3 = tf.add_paragraph()
p3.text = "Operational Strategy & Implementation Architecture"
p3.font.name = "Arial"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

# ==========================================
# SLIDE 2: MANUAL BOTTLENECK VS. INTELLIGENT LEANING
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
apply_background(slide2)
create_title(slide2, "The Paradigm Shift: Manual Bottleneck vs. MetaVal")

card_l = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
card_l.fill.solid()
card_l.fill.fore_color.rgb = CARD_BG
card_l.line.color.rgb = ACCENT_AMBER
card_l.line.width = Pt(1.5)

tf_l = card_l.text_frame
tf_l.word_wrap = True
tf_l.margin_left = Inches(0.3)
tf_l.margin_top = Inches(0.3)

p = tf_l.paragraphs[0]
p.text = "CURRENT MANUAL PROCEDURE"
p.font.bold = True
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_AMBER
p.space_after = Pt(14)

manual_pts = [
    "Fragmented Inputs: Operators manually download varying supplier quotes, invoices, and spec layouts from corporate mailboxes.",
    "Data Loss Vulnerabilities: Physical typing of multi-page tabular data leads to keying fatigue, typos, and format errors.",
    "Operational Delays: Sifting through 10-to-500 page logs creates extreme bottlenecks before procurement choices can be processed."
]
for pt in manual_pts:
    p = tf_l.add_paragraph()
    p.text = "• " + pt
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(10)

card_r = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.93), Inches(1.8), Inches(5.8), Inches(4.8))
card_r.fill.solid()
card_r.fill.fore_color.rgb = CARD_BG
card_r.line.color.rgb = ACCENT_GREEN
card_r.line.width = Pt(1.5)

tf_r = card_r.text_frame
tf_r.word_wrap = True
tf_r.margin_left = Inches(0.3)
tf_r.margin_top = Inches(0.3)

p = tf_r.paragraphs[0]
p.text = "METAVAL AUTOMATION SYSTEM"
p.font.bold = True
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_GREEN
p.space_after = Pt(14)

meta_pts = [
    "Lean Processing: Scrapes incoming emails and attachments hands-free, routing them instantly into a data parser.",
    "Unified Ingestion Matrix: Automatically separates, structuralizes, and standardizes digital and scanned paper layers.",
    "Near-Zero Latency: Translates raw unstructured grids into a pristine data store within moments, keeping procurement agile."
]
for pt in meta_pts:
    p = tf_r.add_paragraph()
    p.text = "• " + pt
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(10)

# ==========================================
# SLIDE 3: THE LAYMAN FLOWCHART
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
apply_background(slide3)
create_title(slide3, "System Execution Flowchart: Inbound File to Database Target")

flow_steps = [
    {"num": "1. INGEST", "title": "Data Stream Caught", "body": "PDF files arrive automatically via email scraping or operator drag-and-drop on the system dashboard."},
    {"num": "2. IDENTIFY", "title": "Extraction Routine", "body": "System detects text blocks. Digital data runs via text extractors; scanned images loop to tesseract OCR page-by-page."},
    {"num": "3. NORMALIZE", "title": "Fuzzy Schema Map", "body": "Raw columns pass through a translation mapping layer, matching supplier variants to fixed company data fields."},
    {"num": "4. EXPORT", "title": "Clean Data Delivery", "body": "Data maps instantly into Postgres JSONB records, exporting flawlessly as structured sheets with one click."}
]

f_width = Inches(2.7)
f_gap = Inches(0.3)
f_start = Inches(0.6)

for i, step in enumerate(flow_steps):
    c_left = f_start + i * (f_width + f_gap)
    box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, Inches(2.2), f_width, Inches(4.2))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = TEXT_MUTED
    
    tf_b = box.text_frame
    tf_b.word_wrap = True
    tf_b.margin_top = Inches(0.2)
    tf_b.margin_left = Inches(0.2)
    tf_b.margin_right = Inches(0.2)
    
    p1 = tf_b.paragraphs[0]
    p1.text = step["num"]
    p1.font.bold = True
    p1.font.size = Pt(11)
    p1.font.color.rgb = ACCENT_GREEN
    p1.space_after = Pt(4)
    
    p2 = tf_b.add_paragraph()
    p2.text = step["title"]
    p2.font.bold = True
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_after = Pt(12)
    
    p3 = tf_b.add_paragraph()
    p3.text = step["body"]
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_MUTED

# ==========================================
# SLIDE 4: STRATEGIC BUSINESS ADVANTAGES (PART 1)
# ==========================================
slide4 = prs.slides.add_slide(slide_layout)
apply_background(slide4)
create_title(slide4, "Strategic Business Advantages: Lean & Accountable")

adv_box = slide4.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.13), Inches(5.0))
tf_adv = adv_box.text_frame
tf_adv.word_wrap = True

adv_1_4 = [
    ("1. Lean Process Optimization:", "Eliminates non-value-added clerical touchpoints. The primary goal is leaning out the data processing pipeline, shifting employee bandwidth from data entry to strategic vendor analysis."),
    ("2. Frozen Accountability & Total Traceability:", "Essential for industrial compliance. Locks down exact chronological context on when, how, who, and where component parts were quotationed, routed, or sold across the enterprise."),
    ("3. Defined Process Boundaries:", "Assigns precise ownership vectors to every segment of data integration. Prevents tracking gaps, locks down data validity, and highlights step responsibility."),
    ("4. Eliminating Waste (Non-Productive Action Cleansed):", "Kills repetitive data manipulation, double manual review cycles, and spreadsheet corrections. Removes operational delays entirely.")
]

for title, desc in adv_1_4:
    p = tf_adv.add_paragraph()
    p.text = "📈 " + title + " "
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_GREEN
    
    r = p.add_run()
    r.text = desc
    r.font.bold = False
    r.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(16)

# ==========================================
# SLIDE 5: STRATEGIC BUSINESS ADVANTAGES (PART 2)
# ==========================================
slide5 = prs.slides.add_slide(slide_layout)
apply_background(slide5)
create_title(slide5, "Strategic Business Advantages: Growth & System Integration")

adv_box2 = slide5.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.13), Inches(5.0))
tf_adv2 = adv_box2.text_frame
tf_adv2.word_wrap = True

adv_5_8 = [
    ("5. Driving Scalable Business Volume:", "Faster parsing speeds let engineering teams respond to customer quotes instantly, capturing higher transaction volumes and winning more market opportunities."),
    ("6. Democratized Central Control (Single-Screen Dashboard):", "Consolidates multi-department actions into one single accessible dashboard monitor. Breaks data siloes and provides cross-team insight on any device."),
    ("7. Real-Time Integration with Enterprise Resource Planning (ERP):", "Keeps master databases up to date with automated record pipelines, removing the delayed synchronization lag standard in manual data entry loops."),
    ("8. Data-Driven Insights (Daily Success Metrics):", "The telemetry layer computes and prints daily success rates on the dashboard, giving managers immediate clarity on parsing quality trends.")
]

for title, desc in adv_5_8:
    p = tf_adv2.add_paragraph()
    p.text = "🚀 " + title + " "
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_GREEN
    
    r = p.add_run()
    r.text = desc
    r.font.bold = False
    r.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(16)

# ==========================================
# SLIDE 6: THE INTERACTIVE OPERATIONS DASHBOARD LAYOUT
# ==========================================
slide6 = prs.slides.add_slide(slide_layout)
apply_background(slide6)
create_title(slide6, "The Enterprise Operations Control Screen Layout")

container = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.4))
container.fill.solid()
container.fill.fore_color.rgb = BG_COLOR
container.line.color.rgb = TEXT_MUTED
container.line.width = Pt(1.5)

sidebar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(2.2), Inches(5.4))
sidebar.fill.solid()
sidebar.fill.fore_color.rgb = CARD_BG
sidebar.line.fill.background()
sidebar.text_frame.paragraphs[0].text = " MetaVal OS\n\n 📊 Overview\n 🗃️ Traceability\n ⚙️ ERP Synced\n 💬 Anomalies Form"
sidebar.text_frame.paragraphs[0].font.size = Pt(12)
sidebar.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

titles = ["TOTAL UPLOADS", "EMAIL REVERTS", "DAILY UPDATES", "DAILY SUCCESS RATE"]
vals = ["14,250 Files", "0 Pending ✅", "+342 Records", "98.4% Accuracy Ticker"]
colors = [TEXT_LIGHT, ACCENT_GREEN, TEXT_LIGHT, ACCENT_GREEN]

for i in range(4):
    bx = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0 + i*2.35), Inches(1.7), Inches(2.2), Inches(0.9))
    bx.fill.solid()
    bx.fill.fore_color.rgb = CARD_BG
    bx.line.fill.background()
    
    tf_k = bx.text_frame
    tf_k.paragraphs[0].text = titles[i]
    tf_k.paragraphs[0].font.size = Pt(8)
    tf_k.paragraphs[0].font.color.rgb = TEXT_MUTED
    
    p2 = tf_k.add_paragraph()
    p2.text = vals[i]
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = colors[i]

act_bar1 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(2.8), Inches(4.5), Inches(0.8))
act_bar1.fill.solid()
act_bar1.fill.fore_color.rgb = CARD_BG
act_bar1.line.color.rgb = TEXT_MUTED
act_bar1.text_frame.paragraphs[0].text = "📤 Manual Ingestion Box (Multer Hook)"
act_bar1.text_frame.paragraphs[0].font.size = Pt(11)
act_bar1.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT

act_bar2 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(2.8), Inches(4.7), Inches(0.8))
act_bar2.fill.solid()
act_bar2.fill.fore_color.rgb = CARD_BG
act_bar2.line.color.rgb = ACCENT_GREEN
act_bar2.text_frame.paragraphs[0].text = "📥 Direct Excel Output Export (.xlsx)"
act_bar2.text_frame.paragraphs[0].font.size = Pt(11)
act_bar2.text_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN

grid_m = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(3.8), Inches(9.4), Inches(2.8))
grid_m.fill.solid()
grid_m.fill.fore_color.rgb = CARD_BG
grid_m.line.fill.background()

tf_m = grid_m.text_frame
tf_m.word_wrap = True
p = tf_m.paragraphs[0]
p.text = "  Traceability Ledger — Real-Time Parts Auditing Logs"
p.font.bold = True
p.font.size = Pt(11)
p.font.color.rgb = TEXT_LIGHT
p.space_after = Pt(8)

p_h = tf_m.add_paragraph()
p_h.text = "   Part Serial      Quoted To     Status            Accountability Signature"
p_h.font.size = Pt(10)
p_h.font.name = "Consolas"
p_h.font.color.rgb = TEXT_MUTED
p_h.space_after = Pt(4)

p_r = tf_m.add_paragraph()
p_r.text = "   P-100-9921     Logistics B    [VALIDATED ✅]     Operator ID: #9042 (Synced to ERP)"
p_r.font.size = Pt(10)
p_r.font.name = "Consolas"
p_r.font.color.rgb = ACCENT_GREEN

# ==========================================
# SLIDE 7: ACTIVE PRODUCTION TECH STACK
# ==========================================
slide7 = prs.slides.add_slide(slide_layout)
apply_background(slide7)
create_title(slide7, "The Production Technology Stack & Dependencies")

techs = [
    {"title": "Core Runtime", "tools": "Node.js (Express) & React.js", "desc": "Forms a clean, full-stack environment running VITE on the frontend for snappy UI state changes and Express for file routing."},
    {"title": "Extraction Engines", "tools": "pdf-parse / Tesseract.js", "desc": "Uses dual-path parsing logic: instant text reads for selectable files, and automatic pdf-poppler rendering for image-scanned OCR files."},
    {"title": "Grid Reconstruction", "tools": "pdf2json matrix maps", "desc": "Extracts underlying layout coordinates to reconstruct cells into precise JSON structures, preserving row/column relationships."},
    {"title": "Enterprise Database", "tools": "PostgreSQL (JSONB Storage)", "desc": "Saves varying tabular data inside high-performance JSONB formats, allowing deep index searches without altering static core database systems."},
    {"title": "Output System", "tools": "ExcelJS Engine", "desc": "Bypasses risky, format-breaking CSV translation entirely, exporting records straight from memory buffers into fixed cell locations."}
]

for i, tech in enumerate(techs):
    t_top = 1.6 + (i * 1.1)
    
    badge = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(t_top), Inches(2.8), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = CARD_BG
    badge.line.color.rgb = ACCENT_GREEN
    
    tf_bd = badge.text_frame
    tf_bd.word_wrap = True
    p = tf_bd.paragraphs[0]
    p.text = tech["title"]
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_GREEN
    
    desc_box = slide7.shapes.add_textbox(Inches(3.6), Inches(t_top - 0.1), Inches(9.1), Inches(0.8))
    tf_ds = desc_box.text_frame
    tf_ds.word_wrap = True
    
    p1 = tf_ds.paragraphs[0]
    p1.text = tech["tools"]
    p1.font.bold = True
    p1.font.size = Pt(13)
    p1.font.color.rgb = TEXT_LIGHT
    
    p2 = tf_ds.add_paragraph()
    p2.text = tech["desc"]
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED

# Save completed presentation
prs.save("Project_MetaVal_Presentation.pptx")
print("Presentation compiled perfectly as 'Project_MetaVal_Presentation.pptx'!")